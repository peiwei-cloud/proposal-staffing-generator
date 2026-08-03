# -*- coding: utf-8 -*-
"""
標案專案人力與組織圖 Web 產生器
==================================
讀取 Template_Staffing.xlsx（人員資料）與 Photos.zip（大頭照），
一鍵產出：
  Output 1：人員組織架構圖 (.pptx，原生可編輯圖形)
  Output 2：主要工作人力配置表 (.docx 表格 / .pptx 表格)
  Output 3：專案主要人員介紹（組長以上）(.docx)

作者：Claude (Anthropic) — 依需求規格開發
"""

import io
import re
import math
import zipfile
import hashlib
import time
import random
from datetime import datetime

import pandas as pd
import streamlit as st
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from docx import Document
from docx.shared import Cm, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL, WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn as docx_qn
from docx.oxml import OxmlElement

try:
    from google import genai as gemini_genai
    from google.genai import errors as gemini_errors
except ImportError:
    gemini_genai = None
    gemini_errors = None


# =====================================================================
# 常數與樣式設定
# =====================================================================

REQUIRED_COLUMNS = [
    "Layer", "Role", "GroupName", "Name", "Title", "Company", "Badges", "PhotoName",
    "YearsOfExp", "Degree", "JobDescription", "Expertise", "BioNarrative",
]

# Excel 欄位名稱有時會被匯出成含中文說明的樣式，這裡做寬鬆對應
COLUMN_ALIASES = {
    "BioNarrative(經歷與描述）": "BioNarrative",
    "BioNarrative(經歷與描述)": "BioNarrative",
}

# 徽章代碼 -> 全名，可依專案需求自行調整
BADGE_MAP = {
    "技": "技師", "碩": "碩士", "博": "博士", "品": "品質管理人員",
    "安": "勞工安全衛生人員", "採": "採購專業人員", "乙": "乙級技術士",
    "甲": "甲級技術士", "景": "景觀技師", "土": "土木技師", "水": "水利技師",
}

BADGE_COLORS = {
    "技": "1F4E79", "碩": "2E7D32", "博": "6A1B9A", "品": "B8860B",
    "安": "C62828", "採": "00838F", "乙": "455A64", "甲": "37474F",
    "景": "558B2F", "土": "5D4037", "水": "0277BD",
}
DEFAULT_BADGE_COLOR = "607D8B"

# 未知徽章代碼的擴充色盤：以 MD5 雜湊決定索引，確保同一代碼在同一次執行/不同次執行間顏色皆固定不變
FALLBACK_BADGE_PALETTE = [
    "8E44AD", "16A085", "D35400", "2980B9", "C2185B", "6D4C41",
    "00897B", "5E35B1", "F4511E", "3949AB", "7CB342", "AD1457",
    "00695C", "4527A0", "EF6C00", "1565C0", "AA00FF", "2E7D32",
]


def get_badge_color(code: str) -> str:
    """回傳徽章代碼對應的十六進位顏色；未知代碼以確定性 Hash 從擴充色盤中分配固定顏色，不再全部落入死板灰藍色"""
    if code in BADGE_COLORS:
        return BADGE_COLORS[code]
    digest = hashlib.md5(code.encode("utf-8")).hexdigest()
    idx = int(digest, 16) % len(FALLBACK_BADGE_PALETTE)
    return FALLBACK_BADGE_PALETTE[idx]

# 主色系（淺色質感風格）
COLOR_BG = RGBColor(0xF7, 0xF9, 0xFC)
COLOR_NAVY = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_NAVY_LIGHT = RGBColor(0x3E, 0x5C, 0x86)
COLOR_HEADER_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
COLOR_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CARD_BORDER = RGBColor(0xD8, 0xDF, 0xE8)
COLOR_TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
COLOR_TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_FOOTER_BG = RGBColor(0xEE, 0xF1, 0xF6)
COLOR_PLACEHOLDER = RGBColor(0xC9, 0xCF, 0xD8)

# 全域中文字型（PPTX / DOCX 皆套用，含 EastAsia 屬性設定）
GLOBAL_FONT_NAME = "Microsoft JhengHei"


def classify_leadership_zone(row) -> str:
    """依 Layer / Role 規則將人員分類至左/中(subtop)/右三區
    優先規則（Layer 欄位明確指定時強制生效，不受 Role 關鍵字影響）：
      - Layer=='Advisor' → 強制左區
      - Layer=='SubTop'  → 強制歸入中欄的 subtop 堆疊
    其餘依 Role 關鍵字判斷：
      左區：Role 含「顧問」「品質督導」
      中區：Layer=='Top' 且 Role 含「計畫主持人」（未命中任何規則的 Top 預設歸中區，避免版面遺漏）
      右區：Layer=='Top' 且 Role 含「協同」或「代表廠商」
    """
    role = row.get("Role", "") or ""
    layer = row.get("Layer", "")
    if layer == "Advisor":
        return "left"
    if layer == "SubTop":
        return "subtop"
    if ("顧問" in role) or ("品質督導" in role):
        return "left"
    if layer == "Top" and ("協同" in role or "代表廠商" in role):
        return "right"
    if layer == "Top":
        return "center"  # 含「計畫主持人」或未命中關鍵字之 Top 人員，皆歸中區
    return "left"  # 理論上不會執行到此（僅 Top/Advisor/SubTop 會呼叫本函式）


# =====================================================================
# 資料載入與容錯處理
# =====================================================================

# Layer 欄位標準化對照（不分大小寫、去除空白後比對），避免使用者輸入 top/TOP/subtop 等造成人員消失
LAYER_CANONICAL_MAP = {
    "top": "Top", "subtop": "SubTop", "advisor": "Advisor",
    "middle": "Middle", "groupleader": "GroupLeader", "groupmember": "GroupMember",
    "subcontractor": "Subcontractor", "partner": "Subcontractor",
}


def normalize_layer_value(value: str) -> str:
    """將 Layer 欄位值去除頭尾空白並標準化大小寫；無法辨識的值原樣保留（去空白），不強制變更避免誤判"""
    v = (value or "").strip()
    if not v:
        return v
    key = v.lower().replace(" ", "").replace("_", "")
    return LAYER_CANONICAL_MAP.get(key, v)


def load_staffing_excel(file) -> pd.DataFrame:
    """讀取 Excel 並做欄位對應、缺值容錯"""
    df = pd.read_excel(file, dtype=str)
    df = df.rename(columns={k: v for k, v in COLUMN_ALIASES.items() if k in df.columns})
    # 找不到的欄位一律補空字串，避免程式崩潰
    for col in REQUIRED_COLUMNS:
        if col not in df.columns:
            df[col] = ""
    df = df[REQUIRED_COLUMNS].copy()
    for col in REQUIRED_COLUMNS:
        df[col] = df[col].fillna("").astype(str).str.strip()
    # Layer 欄位大小寫防呆標準化（top/TOP/subtop 等一律轉為標準形式）
    df["Layer"] = df["Layer"].apply(normalize_layer_value)
    # YearsOfExp 轉數值，轉換失敗補 0
    df["YearsOfExp_num"] = pd.to_numeric(df["YearsOfExp"], errors="coerce").fillna(0)
    return df


def load_photos_zip(file) -> dict:
    """讀取 Photos.zip，回傳 {檔名(不分大小寫): bytes}"""
    photos = {}
    try:
        with zipfile.ZipFile(file) as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                base = name.split("/")[-1]
                if not base or base.startswith("."):
                    continue
                try:
                    photos[base.lower()] = zf.read(name)
                except Exception:
                    continue
    except Exception as e:
        st.warning(f"照片 Zip 讀取失敗，將全部以預留方塊代替：{e}")
    return photos


def get_photo_bytes(photos: dict, photo_name: str):
    """依 PhotoName 找對應照片；找不到回傳 None（呼叫端需自行繪製灰色預留方塊）"""
    if not photo_name or photo_name.strip() in ("", "—", "-"):
        return None
    key = photo_name.strip().lower()
    if key in photos:
        return photos[key]
    # 容錯：嘗試補副檔名
    stem = key.rsplit(".", 1)[0]
    for ext in (".jpg", ".jpeg", ".png"):
        if (stem + ext) in photos:
            return photos[stem + ext]
    return None


def safe_open_image(raw_bytes):
    """驗證圖片是否可被 Pillow 開啟，避免匯出時崩潰"""
    if not raw_bytes:
        return None
    try:
        img = Image.open(io.BytesIO(raw_bytes))
        img.verify()
        return raw_bytes
    except Exception:
        return None


# =====================================================================
# 統計數據計算（頁尾時間軸）
# =====================================================================

def compute_footer_stats(df: pd.DataFrame) -> dict:
    avg_years = 0
    masters = 0
    technicians = 0
    groups = 0
    try:
        valid_years = df["YearsOfExp_num"].replace(0, pd.NA).dropna()
        avg_years = int(round(valid_years.mean())) if len(valid_years) else 0
    except Exception:
        avg_years = 0
    try:
        masters = int(df["Degree"].str.contains("碩士|博士", na=False).sum())
    except Exception:
        masters = 0
    try:
        technicians = int(df["Badges"].apply(lambda b: "技" in [x.strip() for x in b.split(",")]).sum())
    except Exception:
        technicians = 0
    try:
        grp_df = df[df["Layer"].isin(["GroupLeader", "GroupMember"])]
        groups = grp_df[grp_df["GroupName"].str.strip().ne("") & grp_df["GroupName"].str.strip().ne("—")]["GroupName"].nunique()
    except Exception:
        groups = 0
    return {
        "avg_years": f"{avg_years} 平均年資",
        "masters": f"{masters} 碩士",
        "technicians": f"{technicians} 技師",
        "groups": f"{groups} 專業分組",
    }


# =====================================================================
# 共用小工具：徽章解析、姓名加註
# =====================================================================

# 全域徽章統一排序規範：未收錄的新增證照代碼，依原出現順序接在最後
BADGE_SORT_ORDER = ["技", "碩", "博", "品", "安", "採", "乙", "甲", "景", "土", "水"]
_BADGE_SORT_INDEX = {b: i for i, b in enumerate(BADGE_SORT_ORDER)}


def sort_badges(badges: list) -> list:
    """依全域徽章排序規範排序；未收錄代碼保持原相對順序，統一接在已收錄徽章之後"""
    return sorted(badges, key=lambda b: _BADGE_SORT_INDEX.get(b, len(BADGE_SORT_ORDER)))


def parse_badges(badge_str: str):
    if not badge_str or badge_str.strip() in ("", "—", "-"):
        return []
    badges = [b.strip() for b in badge_str.split(",") if b.strip()]
    return sort_badges(badges)  # 全站統一排序，確保所有卡片／圖例／表格徽章順序一致


def name_with_cert_suffix(row) -> str:
    """依 Output2 樣式：若具技師徽章，姓名後加註 (技師)"""
    badges = parse_badges(row["Badges"])
    name = row["Name"] or "[資料待補]"
    if "技" in badges:
        return f"{name}(技師)"
    return name


def safe_text(value: str, fallback="[資料待補]") -> str:
    value = (value or "").strip()
    return value if value else fallback


def clean_separator_join(parts, sep="／") -> str:
    """過濾空值，並去除每個片段頭尾多餘的分隔符號／空白後才串接，
    避免因來源欄位本身已含斜線（如 Title='資深協理/'）而產生 '／/' 或結尾孤立分隔符號等問題。"""
    cleaned = []
    for p in parts:
        p = (p or "").strip()
        p = p.strip("／/、,， ").strip()
        if p:
            cleaned.append(p)
    return sep.join(cleaned)


# =====================================================================
# ============  Output 1：人員組織架構圖 (PPTX)  ======================
# =====================================================================

def _set_pptx_run_ea_font(run, font_name=GLOBAL_FONT_NAME):
    """補上 PPTX 東亞字型設定（<a:ea typeface=.../>），確保中文於各檢視器正確套用指定字型"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def _set_shape_fill(shape, color: RGBColor, line_color=None, line_width=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(0.75)
    shape.shadow.inherit = False


def _add_textbox(slide, x, y, w, h, text, size=11, bold=False, color=COLOR_TEXT_DARK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name=GLOBAL_FONT_NAME,
                  wrap=True, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color
        _set_pptx_run_ea_font(run, font_name)
    return tb


def _add_photo_or_placeholder(slide, x, y, size, photo_bytes):
    """加入大頭照；若無照片則畫出可編輯的灰色預留方塊"""
    photo_bytes = safe_open_image(photo_bytes)
    if photo_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(photo_bytes), x, y, width=size, height=size)
            return
        except Exception:
            pass
    # 灰色原生預留方塊（可直接在 PowerPoint 中以「變更圖片」取代）
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    _set_shape_fill(ph, COLOR_PLACEHOLDER, line_color=RGBColor(0xAA, 0xAA, 0xAA), line_width=Pt(0.75))
    ph.text_frame.word_wrap = True
    ph.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ph.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "照片\n待補"
    run.font.size = Pt(9)
    run.font.name = GLOBAL_FONT_NAME
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    _set_pptx_run_ea_font(run)


def _add_badge_row(slide, x, y, badges, chip_w=Inches(0.28), chip_h=Inches(0.24), gap=Inches(0.05)):
    cx = x
    for b in badges:
        color = RGBColor.from_string(get_badge_color(b))
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, chip_w, chip_h)
        chip.adjustments[0] = 0.35
        _set_shape_fill(chip, color)
        chip.text_frame.word_wrap = False
        chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        chip.text_frame.margin_left = 0
        chip.text_frame.margin_right = 0
        p = chip.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = b
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.name = GLOBAL_FONT_NAME
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_pptx_run_ea_font(run)
        cx += chip_w + gap
    return cx


def get_used_badges(df: pd.DataFrame) -> list:
    """掃描整份 Excel（所有人員，不分 Layer）實際出現過的徽章代碼，依全域排序規範回傳。
    圖例僅根據此清單繪製，避免顯示本案未使用的預設徽章；碩/博等一律照常列出，不做排除。"""
    seen_set = set()
    seen = []
    for b in df["Badges"]:
        for code in parse_badges(b):
            if code not in seen_set:
                seen_set.add(code)
                seen.append(code)
    return sort_badges(seen)


_LEGEND_ICON = Inches(0.24)
_LEGEND_ROW_H = Inches(0.3)
_LEGEND_COL_W = Inches(1.95)
_LEGEND_TOP_PAD = Inches(0.15)
_LEGEND_BOTTOM_PAD = Inches(0.1)


def compute_legend_height(n_badges: int) -> Emu:
    """試算圖例區塊總高度（含上下留白），供版面配置預先計算使用"""
    if n_badges <= 0:
        return Inches(0)
    rows = math.ceil(n_badges / 2)
    return _LEGEND_TOP_PAD + _LEGEND_ROW_H * rows + _LEGEND_BOTTOM_PAD


def _add_legend(slide, x, y, used_badges):
    """繪製左上角圖例：僅列出本案實際出現的徽章，2 欄排列，色塊沿用 get_badge_color() 動態配色；
    找不到全名時直接顯示代碼本身，避免 KeyError。"""
    if not used_badges:
        return
    for i, code in enumerate(used_badges):
        r, c = divmod(i, 2)
        cx = x + c * _LEGEND_COL_W
        cy = y + _LEGEND_TOP_PAD + r * _LEGEND_ROW_H
        color = RGBColor.from_string(get_badge_color(code))
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, _LEGEND_ICON, _LEGEND_ICON)
        chip.adjustments[0] = 0.25
        _set_shape_fill(chip, color)
        chip.text_frame.word_wrap = False
        chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        chip.text_frame.margin_left = 0
        chip.text_frame.margin_right = 0
        p = chip.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = code
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.name = GLOBAL_FONT_NAME
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_pptx_run_ea_font(run)

        label = BADGE_MAP.get(code, code)  # 圖例找不到全名時，直接顯示代碼本身，避免報錯
        _add_textbox(slide, cx + _LEGEND_ICON + Inches(0.08), cy - Inches(0.01),
                     _LEGEND_COL_W - _LEGEND_ICON - Inches(0.12), _LEGEND_ICON + Inches(0.04),
                     label, size=9, color=COLOR_TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


_ORG_BANNER_H = Inches(0.3)
_SUMMARY_CARD_H = Inches(0.85)


def compute_top_section_height(host_org: str, client_org: str, legend_h: Emu) -> Emu:
    """統一計算「機關列 + 圖例/Summary 卡片列」所需總高度，供試算與繪製階段共用，確保一致不重疊"""
    banner_h = _ORG_BANNER_H if (host_org or client_org) else Inches(0)
    banner_gap = Inches(0.06) if banner_h else Inches(0)
    top_row_h = max(legend_h, _SUMMARY_CARD_H)
    return banner_h + banner_gap + top_row_h + Inches(0.14)


def _add_compact_org_banner(slide, slide_w, host_org: str, client_org: str):
    """頂端單行精簡機關列（高約 0.3in），無填寫時不繪製"""
    if not (host_org or client_org):
        return
    parts = []
    if host_org:
        parts.append(f"主辦機關：{host_org}")
    if client_org:
        parts.append(f"委辦機關：{client_org}")
    text = "　｜　".join(parts)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, _ORG_BANNER_H)
    _set_shape_fill(bar, COLOR_NAVY_LIGHT)
    _add_textbox(slide, 0, 0, slide_w, _ORG_BANNER_H, text, size=11, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)


def extract_expertise_keywords(df: pd.DataFrame, max_n: int = 6) -> list:
    """從 Expertise 欄位自動萃取關鍵詞（依 / 、 , ， 拆解），依首次出現順序取前 max_n 個不重複詞"""
    seen_set = set()
    seen = []
    for val in df["Expertise"]:
        val = (val or "").strip()
        if not val:
            continue
        for part in re.split(r"[/、,，]", val):
            part = part.strip()
            if part and part not in seen_set:
                seen_set.add(part)
                seen.append(part)
                if len(seen) >= max_n:
                    return seen
    return seen


def build_summary_stats(df: pd.DataFrame, expertise_override: str = "") -> dict:
    """團隊學經歷 Summary 卡片所需統計數據：總人數／技師級專家人數／專業涵蓋領域文字"""
    total = len(df)
    try:
        technicians = int(df["Badges"].apply(
            lambda b: "技" in [x.strip() for x in b.split(",")] if b else False).sum())
    except Exception:
        technicians = 0
    override = (expertise_override or "").strip()
    if override:
        expertise_str = override
    else:
        kws = extract_expertise_keywords(df, max_n=6)
        expertise_str = "、".join(kws) if kws else "水利工程、水資源規劃、土木防洪、水土保持、環工設計"
    return {"total": total, "technicians": technicians, "expertise": expertise_str}


def _add_summary_card(slide, x, y, w, summary: dict):
    """右上角團隊學經歷 Summary 卡片：左側藍色圓角區顯示總人數，右側文字區顯示技師人數與專長涵蓋"""
    h = _SUMMARY_CARD_H
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    _set_shape_fill(card, COLOR_CARD_BG, line_color=COLOR_NAVY_LIGHT, line_width=Pt(1.25))

    pad = Inches(0.09)
    blue_w = Inches(1.3)
    blue_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + pad, y + pad, blue_w - pad, h - pad * 2)
    blue_box.adjustments[0] = 0.12
    _set_shape_fill(blue_box, COLOR_NAVY)
    blue_box.text_frame.word_wrap = True
    blue_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    blue_box.text_frame.margin_left = 0
    blue_box.text_frame.margin_right = 0
    p1 = blue_box.text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = f"{summary['total']}名"
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.name = GLOBAL_FONT_NAME
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_pptx_run_ea_font(r1)
    p2 = blue_box.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "團隊成員"
    r2.font.size = Pt(10)
    r2.font.bold = True
    r2.font.name = GLOBAL_FONT_NAME
    r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_pptx_run_ea_font(r2)

    text_x = x + blue_w + Inches(0.12)
    text_w = w - blue_w - Inches(0.24)
    line1 = f"配置 {summary['technicians']} 位技師級專家，具專業執照與高階碩博士學歷"
    line2 = f"專業涵蓋：{summary['expertise']}"
    _add_textbox(slide, text_x, y + Inches(0.12), text_w, Inches(0.32), line1,
                 size=10, bold=False, color=COLOR_TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True)
    _add_textbox(slide, text_x, y + Inches(0.46), text_w, Inches(0.34), line2,
                 size=10, bold=True, color=COLOR_NAVY, align=PP_ALIGN.LEFT, wrap=True)


def get_subcontractor_rows(df: pd.DataFrame) -> list:
    """依 Layer=='Subcontractor'（含 Partner 別名，讀取時已標準化）篩選協力廠商項目"""
    return df[df["Layer"] == "Subcontractor"].to_dict("records")


# 協力廠商卡片配色（方案 B：淺色系莫蘭迪風）—— 標頭深藍灰、內容區極淺天藍灰、純白卡片底圖
COLOR_PARTNER_HEADER = RGBColor(0x2C, 0x3E, 0x50)
COLOR_PARTNER_BODY_BG = RGBColor(0xF4, 0xF7, 0xF9)
COLOR_PARTNER_TEXT = RGBColor(0x33, 0x33, 0x33)

_PB_TOP_PAD = Inches(0.08)
_PB_LINE_H = Inches(0.2)
_PB_BOTTOM_PAD = Inches(0.08)
_PB_MIN_BODY_H = Inches(0.3)


def _partner_box_height(n_partners: int, scale: float = 1.0) -> Emu:
    """協力廠商卡片總高度（標頭 + 內容區），與 _add_partner_box 實際繪製共用同一組常數，確保試算與繪製一致"""
    if n_partners <= 0:
        return Inches(0)
    header_h = _LC_HEADER_H * scale
    body_h = (_PB_TOP_PAD + _PB_LINE_H * n_partners + _PB_BOTTOM_PAD) * scale
    body_h = max(body_h, _PB_MIN_BODY_H * scale)
    return int(header_h + body_h)


def _add_partner_box(slide, x, y, w, partner_rows, scale: float = 1.0):
    """繪製協力廠商卡片：外觀比照領導卡片標準（深色圓角標頭 + 實線外框淺色卡片主體），
    與頂層領導卡片陣容視覺一致；寬度／高度依實際筆數動態調整，不佔用固定空間。"""
    n = len(partner_rows)
    if n == 0:
        return
    header_h = int(_LC_HEADER_H * scale)
    h = _partner_box_height(n, scale)

    # 卡片主體：實線外框、淺色底（取消原虛線外框設計）
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_PARTNER_BODY_BG
    card.line.color.rgb = COLOR_CARD_BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # 標頭：與領導卡片同款深色圓角矩形，統一顯示「協力廠商」
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, header_h)
    header.adjustments[0] = 0.5
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PARTNER_HEADER
    header.line.fill.background()
    header.shadow.inherit = False
    header.text_frame.word_wrap = True
    header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = header.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "協力廠商"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.name = GLOBAL_FONT_NAME
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_pptx_run_ea_font(run)

    # 廠商清單：一列一行，靠左置中對齊，字級與領導卡片職稱說明一致（9.5pt）
    line_h = int(_PB_LINE_H * scale)
    content_x = x + int(Inches(0.14) * scale)
    content_w = w - int(Inches(0.28) * scale)
    cy = y + header_h + int(_PB_TOP_PAD * scale)
    for r in partner_rows:
        item = safe_text(r.get("Title", ""), "") or safe_text(r.get("Expertise", ""), "")
        company = safe_text(r.get("Name", ""), "[廠商待補]")
        line = f"{item}：{company}" if item else company
        _add_textbox(slide, content_x, cy, content_w, line_h, line,
                     size=9.5, bold=False, color=COLOR_PARTNER_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cy += line_h



_LC_HEADER_H = Inches(0.24)
_LC_PHOTO_SIZE = Inches(0.85)
_LC_TOP_PAD = Inches(0.05)
_LC_BOTTOM_PAD = Inches(0.05)
_LC_PHOTO_LEFT_PAD = Inches(0.09)
_LC_INFO_GAP = Inches(0.1)
_LC_RIGHT_PAD = Inches(0.1)


def _leader_card_height(scale: float = 1.0) -> Emu:
    """固定卡片高度（橫向版型下，姓名/職稱/公司/徽章皆與照片同高並排，不再隨欄位增減而變動）"""
    return int((_LC_HEADER_H + _LC_TOP_PAD + _LC_PHOTO_SIZE + _LC_BOTTOM_PAD) * scale)


def _add_leader_card(slide, cx, top, row, photos, card_w=Inches(2.8), scale: float = 1.0):
    """橫向領導卡片：頂部職稱標頭列，下方左側照片、右側姓名/職稱與公司/徽章由上至下排列"""
    header_h = int(_LC_HEADER_H * scale)
    photo_size = int(_LC_PHOTO_SIZE * scale)
    top_pad = int(_LC_TOP_PAD * scale)
    bottom_pad = int(_LC_BOTTOM_PAD * scale)
    photo_left_pad = int(_LC_PHOTO_LEFT_PAD * scale)
    info_gap = int(_LC_INFO_GAP * scale)
    right_pad = int(_LC_RIGHT_PAD * scale)
    card_h = header_h + top_pad + photo_size + bottom_pad
    left = int(cx - card_w / 2)

    # 白底卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.adjustments[0] = 0.08
    _set_shape_fill(card, COLOR_CARD_BG, line_color=COLOR_CARD_BORDER, line_width=Pt(1))

    # 深藍色職稱標頭（僅上緣圓角視覺，實務上以矩形疊加圓角卡片頂部即可達成一致外觀）
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, header_h)
    header.adjustments[0] = 0.5
    _set_shape_fill(header, COLOR_NAVY)
    header.text_frame.word_wrap = True
    header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = header.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = safe_text(row["Role"], row["Role"] or "職稱待補")
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.name = GLOBAL_FONT_NAME
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_pptx_run_ea_font(run)

    # 照片（左側，與右側資訊區同一水平帶）
    photo_x = left + photo_left_pad
    photo_y = top + header_h + top_pad
    _add_photo_or_placeholder(slide, photo_x, photo_y, photo_size, get_photo_bytes(photos, row["PhotoName"]))

    # 右側資訊區：姓名 / 職稱與公司 / 徽章，由上至下排列，字級固定 9~12pt 不隨 scale 縮放
    info_x = photo_x + photo_size + info_gap
    info_w = left + card_w - right_pad - info_x
    info_w = max(info_w, Inches(0.6))

    name_h = int(photo_size * 0.30)
    title_h = int(photo_size * 0.24)
    badge_h = int(photo_size * 0.26)

    name_y = photo_y
    _add_textbox(slide, info_x, name_y, info_w, name_h, safe_text(row["Name"]),
                 size=11, bold=True, color=COLOR_TEXT_DARK, align=PP_ALIGN.LEFT)

    company = safe_text(row.get("Company", ""), "")
    title_text = clean_separator_join([safe_text(row["Title"], ""), company], sep="｜")
    title_y = name_y + name_h
    _add_textbox(slide, info_x, title_y, info_w, title_h, title_text,
                 size=9, color=COLOR_TEXT_GRAY, align=PP_ALIGN.LEFT)

    badges = parse_badges(row["Badges"])
    if badges:
        badge_y = title_y + title_h
        chip_h = min(Inches(0.2), badge_h)
        _add_badge_row(slide, info_x, badge_y, badges, chip_w=Inches(0.26), chip_h=chip_h, gap=Inches(0.04))

    return left, top, card_w, card_h


_GB_HEADER_H = Inches(0.24)
_GB_PRE_LEADER_GAP = Inches(0.05)
_GB_LEADER_PHOTO_SIZE = Inches(0.55)
_GB_POST_LEADER_GAP = Inches(0.08)
_GB_DIVIDER_GAP = Inches(0.08)
_GB_MEMBER_LABEL_H = Inches(0.15)
_GB_MEMBER_LINE_H = Inches(0.2)
_GB_BOTTOM_PAD = Inches(0.08)
_GB_MIN_H = Inches(1.4)


def _group_block_height(n_members: int, scale: float = 1.0) -> Emu:
    """與 _add_group_block 實際繪製流程逐段對應的高度公式（同一組具名常數，避免試算與繪製結果不一致）"""
    h = (_GB_HEADER_H + _GB_PRE_LEADER_GAP + _GB_LEADER_PHOTO_SIZE + _GB_POST_LEADER_GAP
         + _GB_DIVIDER_GAP + _GB_MEMBER_LABEL_H
         + _GB_MEMBER_LINE_H * max(n_members, 1) + _GB_BOTTOM_PAD) * scale
    return int(max(h, _GB_MIN_H * scale))


def _add_group_block(slide, left, top, width, height, group_name, leader_row, member_rows, photos, scale: float = 1.0):
    header_h = int(_GB_HEADER_H * scale)
    # 外框
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    outer.adjustments[0] = 0.03
    _set_shape_fill(outer, COLOR_CARD_BG, line_color=COLOR_CARD_BORDER, line_width=Pt(1))

    # 組名深灰標頭
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, header_h)
    header.adjustments[0] = 0.5
    _set_shape_fill(header, COLOR_HEADER_GRAY)
    header.text_frame.word_wrap = True
    header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = header.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = safe_text(group_name, "組別待補")
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.name = GLOBAL_FONT_NAME
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _set_pptx_run_ea_font(run)

    y = top + header_h + int(_GB_PRE_LEADER_GAP * scale)

    # 組長區塊
    if leader_row is not None:
        photo_size = int(_GB_LEADER_PHOTO_SIZE * scale)
        _add_photo_or_placeholder(slide, left + int(Inches(0.1) * scale), y, photo_size,
                                   get_photo_bytes(photos, leader_row["PhotoName"]))
        info_x = left + int(Inches(0.1) * scale) + photo_size + int(Inches(0.08) * scale)
        info_w = width - (info_x - left) - int(Inches(0.08) * scale)
        _add_textbox(slide, info_x, y, info_w, int(Inches(0.15) * scale), "組長",
                     size=9, bold=True, color=COLOR_NAVY, align=PP_ALIGN.LEFT)
        _add_textbox(slide, info_x, y + int(Inches(0.15) * scale), info_w, int(Inches(0.2) * scale),
                     f"{safe_text(leader_row['Name'])}　{safe_text(leader_row['Title'],'')}",
                     size=10, bold=True, color=COLOR_TEXT_DARK, align=PP_ALIGN.LEFT)
        badges = parse_badges(leader_row["Badges"])
        if badges:
            _add_badge_row(slide, info_x, y + int(Inches(0.34) * scale), badges,
                           chip_w=Inches(0.2), chip_h=Inches(0.18), gap=Inches(0.03))
        y += photo_size + int(_GB_POST_LEADER_GAP * scale)
    else:
        _add_textbox(slide, left + int(Inches(0.1) * scale), y, width - int(Inches(0.2) * scale),
                     int(Inches(0.18) * scale), "組長：[資料待補]",
                     size=9, color=COLOR_TEXT_GRAY, align=PP_ALIGN.LEFT)
        y += int(_GB_LEADER_PHOTO_SIZE * scale) + int(_GB_POST_LEADER_GAP * scale)

    # 分隔線（以細長矩形代替裝飾線，符合可編輯需求）
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + int(Inches(0.1) * scale), y,
                                      width - int(Inches(0.2) * scale), Pt(1))
    _set_shape_fill(divider, COLOR_CARD_BORDER)
    y += int(_GB_DIVIDER_GAP * scale)

    _add_textbox(slide, left + int(Inches(0.1) * scale), y, width - int(Inches(0.2) * scale),
                 int(_GB_MEMBER_LABEL_H * scale), "組員", size=9, bold=True, color=COLOR_NAVY, align=PP_ALIGN.LEFT)
    y += int(_GB_MEMBER_LABEL_H * scale)

    if member_rows:
        for m in member_rows:
            badges = parse_badges(m["Badges"])
            line_h = int(_GB_MEMBER_LINE_H * scale)
            _add_textbox(slide, left + int(Inches(0.1) * scale), y, width - int(Inches(1.05) * scale), line_h,
                         f"{safe_text(m['Name'])} {safe_text(m['Title'],'')}",
                         size=9, color=COLOR_TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            if badges:
                _add_badge_row(slide, left + width - int(Inches(0.95) * scale), y + int(Inches(0.01) * scale),
                               badges, chip_w=Inches(0.2), chip_h=Inches(0.18), gap=Inches(0.03))
            y += line_h
    else:
        _add_textbox(slide, left + int(Inches(0.1) * scale), y, width - int(Inches(0.2) * scale),
                     int(Inches(0.18) * scale), "[資料待補]", size=9, color=COLOR_TEXT_GRAY, align=PP_ALIGN.LEFT)


def _leadership_zones(df: pd.DataFrame, swap_lr: bool = False) -> dict:
    """依規則將 Top/Advisor/SubTop 統一透過 classify_leadership_zone 分類為 left/center/right/subtop；Middle 另外歸類。

    頂層動態平衡：
      - 狀況 A（左右候選皆非空）：維持原樣，左右對稱。
      - 狀況 B（左候選為空、右候選非空，如僅有協同主持人無顧問/督導）：自動將右候選整批搬到左邊，
        避免頂層一邊掛卡片、一邊留白。
      - swap_lr=True 時（使用者手動開關「左右對調」）：整批交換左右候選，在狀況 B 校正之後才套用，
        確保任一情境下都不會有一邊空白。
    """
    zones = {"left": [], "center": [], "right": [], "subtop": []}
    for r in df[df["Layer"].isin(["Top", "Advisor", "SubTop"])].to_dict("records"):
        zone = classify_leadership_zone(r)
        zones.setdefault(zone, []).append(r)

    if not zones["left"] and zones["right"]:
        zones["left"], zones["right"] = zones["right"], []

    if swap_lr:
        zones["left"], zones["right"] = zones["right"], zones["left"]

    zones["middle"] = df[df["Layer"] == "Middle"].to_dict("records")
    return zones


def _stack_height(n_cards: int, scale: float = 1.0, card_gap: Emu = None) -> Emu:
    """計算同一欄 n 張卡片垂直堆疊後的總高度（卡片高度固定，僅視安全網層級縮放）"""
    if n_cards <= 0:
        return 0
    gap = card_gap if card_gap is not None else int(Inches(0.14) * scale)
    return _leader_card_height(scale) * n_cards + gap * max(n_cards - 1, 0)


COLOR_CONNECTOR = COLOR_NAVY_LIGHT  # 連接線改用與標頭同色系的深藍灰，投影時更清晰明顯


def connector_width_pt(scale: float = 1.0) -> float:
    """連接線粗細依版面密度（安全網縮放係數）自適應：版面寬鬆時加粗突顯層級，密集時收細避免打架"""
    return max(1.0, min(2.25, 1.75 * scale))


def _add_connector(slide, x1, y1, x2, y2, straight=True, width_pt: float = 1.5):
    """加入連接線；straight=True 用直線（同欄垂直線），否則用直角(elbow)連接線"""
    conn_type = 1 if straight else 2  # MSO_CONNECTOR: 1=STRAIGHT, 2=ELBOW
    connector = slide.shapes.add_connector(conn_type, int(x1), int(y1), int(x2), int(y2))
    connector.line.color.rgb = COLOR_CONNECTOR
    connector.line.width = Pt(width_pt)
    return connector


def _compute_leadership_block_height(zones: dict, scale: float, gap_v: Emu, card_gap: Emu,
                                      n_partners: int = 0) -> Emu:
    """計算領導層區塊（左／中軸鏈／右+協力廠商）由頂端到最底部所需總高度，供試算與繪製階段共用同一公式"""
    card_h = _leader_card_height(scale)
    left_h = _stack_height(len(zones["left"]), scale, card_gap)
    right_h = _stack_height(len(zones["right"]), scale, card_gap)
    if n_partners > 0:
        right_h += gap_v + _partner_box_height(n_partners, scale)
    center_own_h = _stack_height(len(zones["center"]), scale, card_gap) if zones["center"] else 0
    chain_h = center_own_h
    if zones["subtop"]:
        chain_h += gap_v + _stack_height(len(zones["subtop"]), scale, card_gap)
    if zones["middle"]:
        chain_h += gap_v + card_h
    return max(left_h, right_h, chain_h)


def _compute_groups_total_height(df: pd.DataFrame, scale: float, group_gap_y: Emu):
    group_names = []
    for g in df[df["Layer"].isin(["GroupLeader", "GroupMember"])]["GroupName"]:
        g = g.strip()
        if g and g not in group_names:
            group_names.append(g)
    n = max(len(group_names), 1)
    cols = min(4, n) if n > 0 else 1
    row_max_h = {}
    for idx, gname in enumerate(group_names):
        r, _c = divmod(idx, cols)
        n_members = len(df[(df["Layer"] == "GroupMember") & (df["GroupName"] == gname)])
        bh = _group_block_height(n_members, scale)
        row_max_h[r] = max(row_max_h.get(r, 0), bh)
    total = sum(row_max_h.values()) + group_gap_y * max(len(row_max_h) - 1, 0)
    return total, group_names, cols


def _render_column_stack(slide, cx, top_y, rows, card_w, scale, card_gap, photos, conn_width=1.5):
    """垂直堆疊繪製同一欄卡片，相鄰卡片間以直線連接線相連。回傳 (第一張卡片底部y, 最後一張卡片底部y)"""
    card_h = _leader_card_height(scale)
    y = top_y
    first_bottom = None
    last_bottom = None
    for i, r in enumerate(rows):
        _add_leader_card(slide, cx, y, r, photos, card_w=card_w, scale=scale)
        if last_bottom is not None:
            try:
                _add_connector(slide, cx, last_bottom, cx, y, straight=True, width_pt=conn_width)
            except Exception:
                pass
        if first_bottom is None:
            first_bottom = y + card_h
        last_bottom = y + card_h
        y += card_h + card_gap
    return first_bottom, last_bottom


def build_org_chart_pptx(df: pd.DataFrame, stats: dict, photos: dict,
                          host_org: str = "", client_org: str = "",
                          team_expertise_override: str = "", swap_lr: bool = False) -> io.BytesIO:
    # ---------- 投影片尺寸嚴格鎖定標準 16:9 寬螢幕，不再動態加高 ----------
    slide_w = Inches(13.333)
    slide_h = Inches(7.5)

    zones = _leadership_zones(df, swap_lr=swap_lr)
    used_badges = get_used_badges(df)  # 圖例：本案實際出現的徽章全數列出（含碩/博），依全域排序規範排列
    legend_h = compute_legend_height(len(used_badges))
    summary = build_summary_stats(df, team_expertise_override)
    top_section_h = compute_top_section_height(host_org, client_org, legend_h)
    partner_rows = get_subcontractor_rows(df)
    n_partners = len(partner_rows)

    FOOTER_H = Inches(0.75)
    BOTTOM_MARGIN = Inches(0.08)
    GAP_V_NOM, CARD_GAP_NOM, GROUP_GAP_Y_NOM = Inches(0.2), Inches(0.1), Inches(0.15)
    GAP_V_MIN, CARD_GAP_MIN, GROUP_GAP_Y_MIN = Inches(0.1), Inches(0.06), Inches(0.08)
    SCALE_FLOOR = 0.7
    available = slide_h

    def total_height_for(scale, gap_v, card_gap, group_gap_y):
        leadership_h = _compute_leadership_block_height(zones, scale, gap_v, card_gap, n_partners)
        groups_h, _names, _cols = _compute_groups_total_height(df, scale, group_gap_y)
        return top_section_h + leadership_h + gap_v + groups_h + FOOTER_H + BOTTOM_MARGIN

    # ---------- 三層安全網：① 標準間距 → ② 最小間距 → ③ 等比例幾何縮放（字級固定不變） ----------
    total_nom = total_height_for(1.0, GAP_V_NOM, CARD_GAP_NOM, GROUP_GAP_Y_NOM)
    if total_nom <= available:
        scale, gap_v, card_gap, group_gap_y = 1.0, GAP_V_NOM, CARD_GAP_NOM, GROUP_GAP_Y_NOM
    else:
        total_min = total_height_for(1.0, GAP_V_MIN, CARD_GAP_MIN, GROUP_GAP_Y_MIN)
        if total_min <= available:
            scale, gap_v, card_gap, group_gap_y = 1.0, GAP_V_MIN, CARD_GAP_MIN, GROUP_GAP_Y_MIN
        else:
            fixed_part = top_section_h + FOOTER_H + BOTTOM_MARGIN
            compressible_nominal = max(total_min - fixed_part, Inches(0.01))
            available_for_compressible = max(available - fixed_part, Inches(0.5))
            raw_scale = available_for_compressible / compressible_nominal
            scale = max(SCALE_FLOOR, min(1.0, raw_scale))
            gap_v = int(GAP_V_MIN * scale)
            card_gap = int(CARD_GAP_MIN * scale)
            group_gap_y = int(GROUP_GAP_Y_MIN * scale)

    conn_width = connector_width_pt(scale)

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面

    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    _set_shape_fill(bg, COLOR_BG)
    bg.shadow.inherit = False

    # --- 頂部黃金三元素：機關列 / 圖例 / Summary 卡片 ---
    _add_compact_org_banner(slide, slide_w, host_org, client_org)
    banner_h = _ORG_BANNER_H if (host_org or client_org) else Inches(0)
    banner_gap = Inches(0.06) if banner_h else Inches(0)
    top_row_y = banner_h + banner_gap

    _add_legend(slide, Inches(0.4), top_row_y, used_badges)

    summary_w = Inches(6.6)
    summary_x = slide_w - Inches(0.4) - summary_w
    _add_summary_card(slide, summary_x, top_row_y, summary_w, summary)

    leadership_row_y = top_section_h

    # --- 中軸領導層：標準三欄對稱結構（左／中／右） ---
    # 版面寬鬆時（標準間距、分組數不多）將卡片略為縮窄，拉長左右連接線視覺比例；
    # 版面密集時維持較寬卡片以節省空間
    group_count_probe = df[df["Layer"].isin(["GroupLeader", "GroupMember"])]["GroupName"].nunique()
    slot_w = (slide_w - Inches(1.0)) / 3
    left_cx = Inches(0.5) + slot_w * 0.5
    center_cx = slide_w / 2
    right_cx = Inches(0.5) + slot_w * 2.5
    if scale >= 0.999 and group_count_probe <= 4:
        card_w = min(Inches(2.6), slot_w - Inches(0.15))
    else:
        card_w = min(Inches(3.0), slot_w - Inches(0.15))
    card_h = _leader_card_height(scale)

    left_first_bottom, left_last_bottom = _render_column_stack(
        slide, left_cx, leadership_row_y, zones["left"], card_w, scale, card_gap, photos, conn_width)
    center_first_bottom, center_last_bottom = _render_column_stack(
        slide, center_cx, leadership_row_y, zones["center"], card_w, scale, card_gap, photos, conn_width)
    right_first_bottom, right_last_bottom = _render_column_stack(
        slide, right_cx, leadership_row_y, zones["right"], card_w, scale, card_gap, photos, conn_width)

    # --- 協力廠商：固定於右欄卡片下方，寬度與領導卡片一致，長高依實際筆數動態調整 ---
    if partner_rows:
        partner_top = (right_last_bottom if right_last_bottom is not None else leadership_row_y) + gap_v
        partner_left = right_cx - card_w / 2
        _add_partner_box(slide, partner_left, partner_top, card_w, partner_rows, scale)

    # 連接線：左／右卡片直接水平連接至中央「計畫主持人」卡片左／右邊界（不向下連接至計畫經理）
    if zones["center"]:
        mid_y = leadership_row_y + card_h // 2
        center_left_edge = center_cx - card_w // 2
        center_right_edge = center_cx + card_w // 2
        try:
            if zones["left"]:
                left_right_edge = left_cx + card_w // 2
                _add_connector(slide, left_right_edge, mid_y, center_left_edge, mid_y, straight=True, width_pt=conn_width)
            if zones["right"]:
                right_left_edge = right_cx - card_w // 2
                _add_connector(slide, center_right_edge, mid_y, right_left_edge, mid_y, straight=True, width_pt=conn_width)
        except Exception:
            pass

    # --- 中軸鏈：計畫主持人 -> SubTop -> 計畫經理，垂直直線相連 ---
    y_cursor = center_last_bottom if zones["center"] else leadership_row_y
    if zones["subtop"]:
        subtop_top = y_cursor + gap_v
        try:
            _add_connector(slide, center_cx, y_cursor, center_cx, subtop_top, straight=True, width_pt=conn_width)
        except Exception:
            pass
        _st_first, st_last_bottom = _render_column_stack(
            slide, center_cx, subtop_top, zones["subtop"], card_w, scale, card_gap, photos, conn_width)
        y_cursor = st_last_bottom

    middle_bottom = None
    if zones["middle"]:
        middle_top = y_cursor + gap_v
        try:
            _add_connector(slide, center_cx, y_cursor, center_cx, middle_top, straight=True, width_pt=conn_width)
        except Exception:
            pass
        _add_leader_card(slide, center_cx, middle_top, zones["middle"][0], photos, card_w=card_w, scale=scale)
        middle_bottom = middle_top + card_h

    leadership_block_h = _compute_leadership_block_height(zones, scale, gap_v, card_gap, n_partners)
    groups_top = leadership_row_y + leadership_block_h + gap_v

    # --- 工作組別區 ---
    groups_h, group_names, cols = _compute_groups_total_height(df, scale, group_gap_y)
    margin_x = Inches(0.4)
    gap_x = Inches(0.2)
    avail_w = slide_w - margin_x * 2 - gap_x * (cols - 1)
    block_w = avail_w / cols if cols else avail_w

    group_infos = []
    for idx, gname in enumerate(group_names):
        r, c = divmod(idx, cols)
        leader_row = df[(df["Layer"] == "GroupLeader") & (df["GroupName"] == gname)]
        leader_row = leader_row.to_dict("records")[0] if len(leader_row) else None
        member_rows = df[(df["Layer"] == "GroupMember") & (df["GroupName"] == gname)].to_dict("records")
        block_h = _group_block_height(len(member_rows), scale)
        group_infos.append({"r": r, "c": c, "gname": gname, "leader_row": leader_row,
                             "member_rows": member_rows, "block_h": block_h})

    row_max_h = {}
    for info in group_infos:
        row_max_h[info["r"]] = max(row_max_h.get(info["r"], 0), info["block_h"])

    row_y_offset = {}
    cum = 0
    for r in sorted(row_max_h):
        row_y_offset[r] = cum
        cum += row_max_h[r] + group_gap_y

    for info in group_infos:
        bx = margin_x + info["c"] * (block_w + gap_x)
        by = groups_top + row_y_offset[info["r"]]
        _add_group_block(slide, bx, by, block_w, info["block_h"], info["gname"],
                          info["leader_row"], info["member_rows"], photos, scale=scale)

    # 連接線：計畫經理 -> 各組別（僅連接第一列），主幹＋水平匯流排＋垂直分支之樹狀直角連接
    if zones["middle"] and middle_bottom is not None:
        row0_centers = [margin_x + info["c"] * (block_w + gap_x) + block_w / 2
                         for info in group_infos if info["r"] == 0]
        if row0_centers:
            bus_y = groups_top - int(Inches(0.1) * scale)
            try:
                _add_connector(slide, center_cx, middle_bottom, center_cx, bus_y, straight=True, width_pt=conn_width)
                if len(row0_centers) > 1:
                    _add_connector(slide, min(row0_centers), bus_y, max(row0_centers), bus_y, straight=True, width_pt=conn_width)
                for gc in row0_centers:
                    _add_connector(slide, gc, bus_y, gc, groups_top, straight=True, width_pt=conn_width)
            except Exception:
                pass

    footer_top = slide_h - FOOTER_H

    # --- 頁尾統計時間軸（大數字為刻意的視覺焦點設計，經確認不受 9-12pt 文字規範限制） ---
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_top, slide_w, FOOTER_H)
    _set_shape_fill(footer, COLOR_FOOTER_BG)
    stat_items = [stats.get("avg_years", ""), stats.get("masters", ""),
                  stats.get("technicians", ""), stats.get("groups", "")]
    seg_w = slide_w / len(stat_items)
    for i, item in enumerate(stat_items):
        parts = item.split(" ", 1)
        number = parts[0] if parts else ""
        label = parts[1] if len(parts) > 1 else ""
        cx0 = i * seg_w
        _add_textbox(slide, cx0, footer_top + Inches(0.05), seg_w, Inches(0.4), number,
                     size=20, bold=True, color=COLOR_NAVY)
        _add_textbox(slide, cx0, footer_top + Inches(0.42), seg_w, Inches(0.28), label,
                     size=10, color=COLOR_TEXT_GRAY)
        if i > 0:
            divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx0, footer_top + Inches(0.12), Pt(1), Inches(0.5))
            _set_shape_fill(divider, RGBColor(0xC7, 0xCF, 0xDA))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# =====================================================================
# =======  Output 2：主要工作人力配置表（DOCX / PPTX 表格）  ===========
# =====================================================================

LAYER_CATEGORY_MAP_DEFAULT = {
    "Advisor": "計畫顧問",
    "Middle": "計畫經理",
    "GroupLeader": "組長",
    "GroupMember": "組員",
}


def _category_for_row(row) -> str:
    if row["Layer"] == "Top":
        zone = classify_leadership_zone(row)
        if zone == "right":
            return safe_text(row["Role"], "協同主持人")
        return "計畫主持人"
    if row["Layer"] == "SubTop":
        return safe_text(row["Role"], "次頂層人員")
    return LAYER_CATEGORY_MAP_DEFAULT.get(row["Layer"], row["Role"] or "人員")


def build_staffing_table_rows(df: pd.DataFrame):
    """回傳依 Output2 樣式排序、格式化後的 list[dict]，供表格預覽/匯出共用；
    協力廠商（Layer=='Subcontractor'）不計入本表，僅於 Output1 組織圖獨立區塊呈現"""
    order = {"Top": 0, "SubTop": 1, "Advisor": 2, "Middle": 3, "GroupLeader": 4, "GroupMember": 5}
    df2 = df[df["Layer"] != "Subcontractor"].copy()
    df2["_order"] = df2["Layer"].map(order).fillna(9)
    # Top 排序：主持人優先於協同主持人
    df2["_sub"] = df2.apply(lambda r: 0 if (r["Layer"] == "Top" and classify_leadership_zone(r) != "right") else 1, axis=1)
    df2 = df2.sort_values(["_order", "_sub", "GroupName"], kind="stable")

    rows = []
    for _, r in df2.iterrows():
        rows.append({
            "類別": _category_for_row(r),
            "姓名": name_with_cert_suffix(r),
            "職稱": safe_text(r["Title"], ""),
            "最高學歷科系": safe_text(r["Degree"], ""),
            "擬任工作內容": safe_text(r["JobDescription"], ""),
            "相關經歷與專長": safe_text(r["Expertise"], ""),
        })
    return rows


def _set_docx_run_font(run, size_pt=12, bold=None, font_name=GLOBAL_FONT_NAME):
    """設定 DOCX run 字型，含 w:eastAsia 屬性，確保中文正確套用指定字型"""
    run.font.name = font_name
    run.font.size = DocxPt(size_pt)
    if bold is not None:
        run.font.bold = bold
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(docx_qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(docx_qn("w:eastAsia"), font_name)


def _set_docx_default_font(doc, font_name=GLOBAL_FONT_NAME):
    """設定文件預設樣式字型（保險機制，涵蓋未逐一設定的段落）"""
    try:
        style = doc.styles["Normal"]
        style.font.name = font_name
        rPr = style.element.get_or_add_rPr()
        rFonts = rPr.find(docx_qn("w:rFonts"))
        if rFonts is None:
            rFonts = OxmlElement("w:rFonts")
            rPr.append(rFonts)
        rFonts.set(docx_qn("w:eastAsia"), font_name)
    except Exception:
        pass


def build_staffing_table_docx(df: pd.DataFrame) -> io.BytesIO:
    rows = build_staffing_table_rows(df)
    doc = Document()
    _set_docx_default_font(doc)
    section = doc.sections[0]
    section.orientation = 0

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("主要工作人力配置表")
    _set_docx_run_font(run, size_pt=16, bold=True)

    headers = ["類別", "姓名", "職稱", "最高學歷科系", "擬任工作內容", "相關經歷與專長"]
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = ""
        p = hdr_cells[i].paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(h)
        _set_docx_run_font(run, size_pt=12, bold=True)
        _shade_cell(hdr_cells[i], "1F3A5F")
        run.font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)
        hdr_cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

    for r in rows:
        cells = table.add_row().cells
        for i, h in enumerate(headers):
            cells[i].text = ""
            p = cells[i].paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT if i >= 3 else WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(r[h])
            _set_docx_run_font(run, size_pt=12, bold=False)
            cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out


def _shade_cell(cell, hex_color):
    shd = OxmlElement("w:shd")
    shd.set(docx_qn("w:fill"), hex_color)
    cell._tc.get_or_add_tcPr().append(shd)


def build_staffing_table_pptx(df: pd.DataFrame) -> io.BytesIO:
    rows = build_staffing_table_rows(df)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    headers = ["類別", "姓名", "職稱", "最高學歷科系", "擬任工作內容", "相關經歷與專長"]
    col_widths = [1.3, 1.5, 1.3, 2.6, 2.6, 4.0]

    # 每頁最多容納列數（含表頭）
    rows_per_slide = 10
    chunks = [rows[i:i + rows_per_slide] for i in range(0, len(rows), rows_per_slide)] or [[]]

    for chunk in chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = _add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5),
                                  "主要工作人力配置表", size=22, bold=True, color=COLOR_NAVY, align=PP_ALIGN.LEFT)
        n_rows = len(chunk) + 1
        table_shape = slide.shapes.add_table(n_rows, len(headers), Inches(0.4), Inches(0.9),
                                              Inches(12.5), Inches(0.4) * n_rows)
        table = table_shape.table
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_NAVY
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.name = GLOBAL_FONT_NAME
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    _set_pptx_run_ea_font(run)
        for ri, r in enumerate(chunk, start=1):
            for ci, h in enumerate(headers):
                cell = table.cell(ri, ci)
                cell.text = r[h]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9.5)
                        run.font.name = GLOBAL_FONT_NAME
                        _set_pptx_run_ea_font(run)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# =====================================================================
# ===========  Output 3：主要人員經歷敘述（組長以上）(DOCX)  ============
# =====================================================================

SENIOR_LAYERS = ["Top", "SubTop", "Advisor", "Middle", "GroupLeader"]

# Gemini 模型選單（側邊欄下拉選單選項）；「自訂模型」為特殊值，選中後另跳出文字輸入框
GEMINI_MODEL_OPTIONS = ["gemini-3.6-flash", "gemini-3.5-flash", "gemini-2.5-flash",
                         "gemini-flash-latest", "自訂模型"]

# 備援退避鏈：使用者選擇的模型優先，其餘依序作為 404/NOT_FOUND 時的自動備援
GEMINI_FALLBACK_CHAIN = ["gemini-3.6-flash", "gemini-3.5-flash", "gemini-2.5-flash", "gemini-flash-latest"]


def get_secret_gemini_key() -> str:
    """安全讀取 Streamlit Secrets 中的 GEMINI_API_KEY；未設定 secrets.toml 時不拋出例外"""
    try:
        return (st.secrets.get("GEMINI_API_KEY", "") or "").strip()
    except Exception:
        return ""


def get_effective_gemini_key() -> str:
    """依規則取得實際呼叫用金鑰：優先 Streamlit Secrets，其次側邊欄手動輸入值"""
    return get_secret_gemini_key() or st.session_state.get("gemini_api_key", "").strip()


def build_gemini_model_candidates(selected_model: str, custom_model: str) -> list:
    """組成備援模型清單：[使用者選取/自訂模型] + 預設備援鏈，並自動去除重複項（保留原順序）"""
    primary = custom_model.strip() if selected_model == "自訂模型" else selected_model.strip()
    raw_candidates = [primary] + GEMINI_FALLBACK_CHAIN
    seen = set()
    candidates = []
    for m in raw_candidates:
        m = (m or "").strip()
        if m and m not in seen:
            seen.add(m)
            candidates.append(m)
    return candidates


def _is_model_not_found_error(e: Exception) -> bool:
    """判斷例外是否為「模型不存在 / 404」類型錯誤，用以觸發自動退避切換下一個候選模型"""
    if gemini_errors is not None and isinstance(e, gemini_errors.ClientError):
        code = getattr(e, "code", None)
        if code == 404:
            return True
    msg = str(e).upper()
    return "404" in msg or "NOT_FOUND" in msg or "NOT FOUND" in msg


def _is_server_overloaded_error(e: Exception) -> bool:
    """判斷例外是否為「503 / 服務過載」類型的暫時性錯誤，用以觸發指數退避重試（同一模型）"""
    if gemini_errors is not None and hasattr(gemini_errors, "ServerError") and isinstance(e, gemini_errors.ServerError):
        return True
    msg = str(e).upper()
    return "503" in msg or "UNAVAILABLE" in msg or "OVERLOADED" in msg


GEMINI_MAX_RETRIES_PER_MODEL = 3   # 單一模型遇 503 時的最大重試次數
GEMINI_BACKOFF_BASE_SECONDS = 1.0  # 指數退避基礎秒數（1s, 2s, 4s ...）


def refine_bio_with_gemini(api_key: str, name: str, job_description: str, bio_narrative: str,
                            model_candidates: list) -> tuple:
    """呼叫 Gemini (google-genai) 依擬任工作內容，將原始履歷潤飾為 200 字以內備標敘述。
    - 遇 404 / Model Not Found：無感切換至下一個備援模型。
    - 遇 503 / 服務過載：對同一模型以指數退避（1s→2s→4s）重試，重試仍失敗才換下一個備援模型。
    - 其餘錯誤（金鑰無效、額度超限等）：直接中止並拋出，不逐一嘗試其餘模型。
    回傳 (潤飾後文字, 實際成功使用的模型 ID)。
    """
    if gemini_genai is None:
        raise RuntimeError("尚未安裝 google-genai 套件，請確認 requirements.txt 是否包含 google-genai")
    if not model_candidates:
        raise RuntimeError("沒有可用的 Gemini 模型可供呼叫，請確認模型選單設定")

    job_description = job_description.strip() or "（未提供擬任工作內容）"
    bio_narrative = bio_narrative.strip() or "（未提供原始履歷）"

    prompt = (
        "你是一位政府水利工程標案專家，"
        f"請依據該人員（{name}）的【擬任工作】：{job_description}，"
        f"將其【原始履歷】：{bio_narrative}，"
        "潤飾精簡為一段200字以內、極具評審說服力且專業精準的水利工程備標履歷敘述。"
        "請只輸出潤飾後的履歷內文本身，不要加上任何前綴、標題、項目符號或引號。"
    )

    client = gemini_genai.Client(api_key=api_key)
    last_error = None
    for model_id in model_candidates:
        for attempt in range(GEMINI_MAX_RETRIES_PER_MODEL):
            try:
                response = client.models.generate_content(model=model_id, contents=prompt)
                text = (getattr(response, "text", "") or "").strip()
                if not text:
                    raise RuntimeError("Gemini 未回傳有效內容")
                return text, model_id
            except Exception as e:
                if _is_server_overloaded_error(e):
                    last_error = e
                    if attempt < GEMINI_MAX_RETRIES_PER_MODEL - 1:
                        sleep_s = GEMINI_BACKOFF_BASE_SECONDS * (2 ** attempt) + random.uniform(0, 0.5)
                        time.sleep(sleep_s)
                        continue  # 503：同一模型指數退避後重試
                    break  # 同一模型重試次數用盡，改嘗試下一個備援模型
                if _is_model_not_found_error(e):
                    last_error = e
                    break  # 404 / 模型不存在：無感切換下一個備援模型
                raise  # 非 404/503 錯誤（金鑰無效、額度超限等）：立即中止，不逐一嘗試其餘模型

    raise RuntimeError(f"候選模型（{', '.join(model_candidates)}）皆呼叫失敗，最後錯誤：{last_error}")


def build_bio_docx(df: pd.DataFrame, photos: dict) -> io.BytesIO:
    order = {"Top": 0, "SubTop": 1, "Advisor": 2, "Middle": 3, "GroupLeader": 4}
    df2 = df[df["Layer"].isin(SENIOR_LAYERS)].copy()
    df2["_order"] = df2["Layer"].map(order).fillna(9)
    df2["_sub"] = df2.apply(lambda r: 0 if (r["Layer"] == "Top" and classify_leadership_zone(r) != "right") else 1, axis=1)
    df2 = df2.sort_values(["_order", "_sub"], kind="stable")

    doc = Document()
    _set_docx_default_font(doc)
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("主要人員經歷")
    _set_docx_run_font(run, size_pt=16, bold=True)

    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.columns[0].width = Cm(4.2)
    table.columns[1].width = Cm(10.5)

    hdr = table.rows[0].cells
    hdr[0].text = ""
    _set_docx_run_font(hdr[0].paragraphs[0].add_run("姓名職稱"), size_pt=12, bold=True)
    hdr[1].text = ""
    _set_docx_run_font(hdr[1].paragraphs[0].add_run("專長說明"), size_pt=12, bold=True)
    for c in hdr:
        _shade_cell(c, "1F3A5F")
        for p in c.paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for r in p.runs:
                r.font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)

    for _, r in df2.iterrows():
        cells = table.add_row().cells
        photo_cell, text_cell = cells[0], cells[1]
        photo_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        photo_p = photo_cell.paragraphs[0]
        photo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        raw = safe_open_image(get_photo_bytes(photos, r["PhotoName"]))
        if raw:
            try:
                run = photo_p.add_run()
                run.add_picture(io.BytesIO(raw), width=Cm(3.2))
            except Exception:
                raw = None
        if not raw:
            # 原生可編輯灰色預留方塊（以有底色表格取代圖片，避免崩潰）
            ph_table = photo_cell.add_table(rows=1, cols=1)
            ph_cell = ph_table.rows[0].cells[0]
            _shade_cell(ph_cell, "C9CFD8")
            ph_cell.text = ""
            _set_docx_run_font(ph_cell.paragraphs[0].add_run("照片待補"), size_pt=12, bold=False)
            ph_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

        badges = parse_badges(r["Badges"])
        cert_names = [BADGE_MAP.get(b, b) for b in badges]
        tail = clean_separator_join([safe_text(r["Title"], "")] + cert_names, sep="／")

        role_line = f"{safe_text(r['Role'])}　{safe_text(r['Name'])}"
        if tail:
            role_line += f"　{tail}"

        p1 = text_cell.paragraphs[0]
        p1.text = ""
        run1 = p1.add_run(role_line)
        _set_docx_run_font(run1, size_pt=12, bold=True)

        p2 = text_cell.add_paragraph()
        run2 = p2.add_run(safe_text(r["BioNarrative"]))
        _set_docx_run_font(run2, size_pt=12, bold=False)

    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out


# =====================================================================
# ============================  Streamlit UI  ==========================
# =====================================================================

# =====================================================================
# ============================  Streamlit UI  ==========================
# =====================================================================

st.set_page_config(page_title="標案專案人力與組織圖產生器", layout="wide", page_icon="📋")

# ---------------- 全域 CSS 美化 ----------------
st.markdown("""
<style>
/* 全站質感淺藍背景 */
.stApp { background-color: #FAFCFF; }

/* 頂部標題 Banner：淺藍漸層圓角卡片 */
.banner-box {
    background: linear-gradient(135deg, #EAF4FD 0%, #F7FBFF 100%);
    border: 1px solid #DCE9F7;
    border-radius: 18px;
    padding: 22px 28px;
    margin-bottom: 14px;
}
.banner-box h1 { margin: 0 0 4px 0; font-size: 28px; color: #123A63; }
.banner-box p { margin: 0; color: #4A5A6A; font-size: 14px; }

/* 提示訊息：淺黃質感提醒條 */
.hint-box {
    background: #FFF9E6;
    border: 1px solid #F3E3A0;
    border-radius: 12px;
    padding: 14px 18px;
    color: #6B5B18;
    font-size: 14px;
    margin: 6px 0 16px 0;
}

/* 主題藍：下載按鈕與一般按鈕 */
div[data-testid="stDownloadButton"] button, div[data-testid="stBaseButton-secondary"] {
    background-color: #007AC3 !important;
    color: #fff !important;
    border: none !important;
    border-radius: 10px !important;
    font-weight: 700 !important;
    padding: 0.6em 1em !important;
}
div[data-testid="stDownloadButton"] button:hover { background-color: #005F99 !important; }
div[data-testid="stButton"] button {
    background-color: #EAF4FD !important;
    color: #123A63 !important;
    border: 1px solid #B9D8F0 !important;
    border-radius: 10px !important;
    font-weight: 600 !important;
}

/* 側邊欄：縮小上傳區塊間距 */
section[data-testid="stSidebar"] .block-container { padding-top: 1.2rem; }
section[data-testid="stSidebar"] h3 { margin-bottom: 4px; }
section[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] { padding: 0.6rem; }
section[data-testid="stSidebar"] .stFileUploader { margin-bottom: 0.3rem; }

/* Expander 標頭美化 */
details summary { font-weight: 700; color: #123A63; }
</style>
""", unsafe_allow_html=True)

# ---------------- Session State 初始化 ----------------
if "df" not in st.session_state:
    st.session_state.df = None
if "photos" not in st.session_state:
    st.session_state.photos = {}
if "stats" not in st.session_state:
    st.session_state.stats = {"avg_years": "0 平均年資", "masters": "0 碩士",
                               "technicians": "0 技師", "groups": "0 專業分組"}
if "host_org" not in st.session_state:
    st.session_state.host_org = ""
if "client_org" not in st.session_state:
    st.session_state.client_org = ""
if "gemini_api_key" not in st.session_state:
    st.session_state.gemini_api_key = get_secret_gemini_key()
if "gemini_model_choice" not in st.session_state:
    st.session_state.gemini_model_choice = GEMINI_MODEL_OPTIONS[0]
if "gemini_custom_model" not in st.session_state:
    st.session_state.gemini_custom_model = ""
if "team_expertise_override" not in st.session_state:
    st.session_state.team_expertise_override = ""
if "swap_lr" not in st.session_state:
    st.session_state.swap_lr = False

# ---------------- ① 側邊欄：資料上傳 + Gemini AI 設定 ----------------
with st.sidebar:
    st.markdown("### ① 資料上傳")

    excel_file = st.file_uploader("Template_Staffing.xlsx", type=["xlsx"], label_visibility="visible")
    if excel_file is not None:
        try:
            st.session_state.df = load_staffing_excel(excel_file)
            st.session_state.stats = compute_footer_stats(st.session_state.df)  # 上傳後自動計算，無需點按鈕
            st.success(f"已讀取 {len(st.session_state.df)} 筆人員資料", icon="✅")
        except Exception as e:
            st.error(f"Excel 讀取失敗：{e}")

    zip_file = st.file_uploader("Photos.zip（大頭照）", type=["zip"], label_visibility="visible")
    if zip_file is not None:
        st.session_state.photos = load_photos_zip(zip_file)
        st.success(f"已讀取 {len(st.session_state.photos)} 張照片", icon="✅")

    st.markdown("### ② Gemini AI 設定")
    _secret_key = get_secret_gemini_key()
    if _secret_key:
        st.caption("🔒 已從 Streamlit Secrets 自動讀取 GEMINI_API_KEY")
    st.session_state.gemini_api_key = st.text_input(
        "Gemini API Key", type="password",
        value=st.session_state.gemini_api_key or _secret_key,
        help="若已於 Streamlit Secrets 設定 GEMINI_API_KEY，將自動帶入並優先使用；"
             "此欄位僅作為備用手動輸入，金鑰僅存於本次瀏覽器工作階段，不會被儲存。",
    )
    st.session_state.gemini_model_choice = st.selectbox(
        "Gemini 模型選取", GEMINI_MODEL_OPTIONS,
        index=GEMINI_MODEL_OPTIONS.index(st.session_state.gemini_model_choice)
        if st.session_state.gemini_model_choice in GEMINI_MODEL_OPTIONS else 0,
        help="潤飾時若所選模型回傳 404 / Model Not Found，將自動無感切換至下一個備援模型。",
    )
    if st.session_state.gemini_model_choice == "自訂模型":
        st.session_state.gemini_custom_model = st.text_input(
            "自訂 Model ID", value=st.session_state.gemini_custom_model,
            placeholder="例如：gemini-2.5-pro",
        )

# ---------------- 頂部標題 Banner ----------------
st.markdown("""
<div class="banner-box">
    <h1>📋 標案專案人力與組織圖 Web 產生器</h1>
    <p>上傳人員資料 Excel 與照片 Zip，即可預覽並一鍵匯出組織圖、人力配置表、人員經歷說明書</p>
</div>
""", unsafe_allow_html=True)

df = st.session_state.df

if df is None:
    st.markdown('<div class="hint-box">👉 請先於左側上傳人員資料 Excel（Template_Staffing.xlsx 格式）</div>',
                unsafe_allow_html=True)
    st.stop()

# ---------------- ② 標案參數設定面板（可折疊） ----------------
with st.expander("⚙️ 標案參數設定面板", expanded=False):
    col_org, col_stats = st.columns(2)

    with col_org:
        st.markdown("**機關名稱**")
        st.session_state.host_org = st.text_input("主辦機關", value=st.session_state.host_org)
        st.session_state.client_org = st.text_input("委辦機關", value=st.session_state.client_org)
        st.markdown("**團隊 Summary 卡片**")
        st.session_state.team_expertise_override = st.text_input(
            "團隊專業涵蓋領域（選填）", value=st.session_state.team_expertise_override,
            placeholder="留空則自動從 Expertise 欄位萃取前 6 個關鍵詞",
            help="填寫後將取代自動萃取結果，直接顯示於組織圖右上角 Summary 卡片。",
        )
        st.markdown("**頂層版面**")
        st.session_state.swap_lr = st.checkbox(
            "左右對調（協同主持人置左、計畫顧問置右）",
            value=st.session_state.swap_lr,
            help="預設左區為計畫顧問/品質督導、右區為協同主持人/代表廠商；勾選後整批對調。",
        )

    with col_stats:
        st.markdown("**頁尾數據時間軸微調**")
        st.caption("Excel 上傳成功後已自動計算，如需重算或還原可點擊下方按鈕。")
        if st.button("🔄 重新計算", use_container_width=True):
            st.session_state.stats = compute_footer_stats(st.session_state.df)

        s1, s2 = st.columns(2)
        with s1:
            st.session_state.stats["avg_years"] = st.text_input("平均年資", st.session_state.stats["avg_years"])
            st.session_state.stats["technicians"] = st.text_input("技師執照數", st.session_state.stats["technicians"])
        with s2:
            st.session_state.stats["masters"] = st.text_input("碩博士人數", st.session_state.stats["masters"])
            st.session_state.stats["groups"] = st.text_input("專業分組數", st.session_state.stats["groups"])

host_org = st.session_state.host_org
client_org = st.session_state.client_org

# ---------------- ③ 匯出下載專區（主畫面頂部，三欄併排） ----------------
st.markdown("#### 📤 匯出下載專區")
exp_col1, exp_col2, exp_col3 = st.columns(3)

with exp_col1:
    try:
        pptx1 = build_org_chart_pptx(df, st.session_state.stats, st.session_state.photos,
                                      host_org, client_org, st.session_state.team_expertise_override,
                                      st.session_state.swap_lr)
        st.download_button("⬇️ Output 1　組織架構圖 (.pptx)", data=pptx1,
                            file_name="人員組織架構圖.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True)
    except Exception as e:
        st.error(f"Output 1 產生失敗：{e}")

with exp_col2:
    try:
        docx2 = build_staffing_table_docx(df)
        st.download_button("⬇️ Output 2　人力配置表 (.docx)", data=docx2,
                            file_name="主要工作人力配置表.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True)
    except Exception as e:
        st.error(f"Output 2 產生失敗：{e}")

with exp_col3:
    try:
        docx3 = build_bio_docx(df, st.session_state.photos)
        st.download_button("⬇️ Output 3　人員經歷說明書 (.docx)", data=docx3,
                            file_name="專案主要人員介紹.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True)
    except Exception as e:
        st.error(f"Output 3 產生失敗：{e}")

st.markdown("")

# ---------------- ④ 下方預覽區：Tab 1~3 ----------------
tab1, tab2, tab3 = st.tabs(["🗂️ 專案組織圖預覽 (Output 1)", "📊 人力配置表預覽 (Output 2)", "🧑‍💼 人員經歷敘述預覽 (Output 3)"])

# ---- Tab 1: 組織圖 HTML 概覽預覽 ----
with tab1:
    st.caption("以下為結構化概覽預覽，實際版面、字型與間距請以上方下載之 .pptx 檔案為準。")

    def badge_html(badges):
        chips = []
        for b in badges:
            color = "#" + get_badge_color(b)
            chips.append(f"<span style='background:{color};color:#fff;border-radius:6px;"
                         f"padding:1px 6px;font-size:11px;margin-right:3px;'>{b}</span>")
        return "".join(chips)

    def person_card(row, small=False):
        badges = parse_badges(row["Badges"])
        photo = get_photo_bytes(st.session_state.photos, row["PhotoName"])
        photo_html = ""
        if safe_open_image(photo):
            import base64
            b64 = base64.b64encode(photo).decode()
            photo_html = f"<img src='data:image/png;base64,{b64}' style='width:56px;height:56px;object-fit:cover;border-radius:6px;'/>"
        else:
            photo_html = "<div style='width:56px;height:56px;background:#C9CFD8;border-radius:6px;" \
                         "display:flex;align-items:center;justify-content:center;font-size:10px;color:#555;'>照片待補</div>"
        company = safe_text(row.get("Company", ""), "")
        company_html = f"<div style='font-size:10px;color:#888;'>{company}</div>" if company else ""
        # 注意：整段回傳字串刻意不使用多行縮排的三引號字串，
        # 因為 Streamlit markdown 若偵測到 4 個以上前導空白會誤判為程式碼區塊，導致 HTML 原始碼亂碼顯示而非正確渲染。
        return (
            "<div style='display:flex;gap:8px;align-items:center;background:#fff;border:1px solid #D8DFE8;"
            "border-radius:8px;padding:6px 8px;margin-bottom:6px;'>"
            f"{photo_html}"
            "<div>"
            f"<div style='font-weight:700;font-size:13px;'>{safe_text(row['Name'])}</div>"
            f"<div style='font-size:11px;color:#666;'>{safe_text(row['Title'],'')}</div>"
            f"{company_html}"
            f"<div style='margin-top:2px;'>{badge_html(badges)}</div>"
            "</div>"
            "</div>"
        )

    zones = _leadership_zones(df, swap_lr=st.session_state.swap_lr)
    center_stack = zones["center"] + zones["subtop"]
    middles = zones["middle"]
    partner_rows_preview = get_subcontractor_rows(df)

    org_header = ""
    if host_org or client_org:
        org_header = "<div style='text-align:center;margin-bottom:10px;'>"
        if host_org:
            org_header += f"<span style='background:#3E5C86;color:#fff;padding:4px 14px;border-radius:14px;margin-right:8px;font-size:13px;'>主辦機關：{host_org}</span>"
        if client_org:
            org_header += f"<span style='background:#3E5C86;color:#fff;padding:4px 14px;border-radius:14px;font-size:13px;'>委辦機關：{client_org}</span>"
        org_header += "</div>"
    st.markdown(org_header, unsafe_allow_html=True)

    st.caption("標準三欄對稱結構：左區（顧問／品質督導）｜中區（計畫主持人＋次頂層 SubTop）｜右區（協同／代表廠商）")
    zone_l, zone_c, zone_r = st.columns(3)
    zone_titles = {"left": "左區", "center": "中區", "right": "右區"}
    for zone_key, col, people in [("left", zone_l, zones["left"]),
                                   ("center", zone_c, center_stack),
                                   ("right", zone_r, zones["right"])]:
        with col:
            st.markdown(f"<div style='text-align:center;color:#8794A6;font-size:11px;font-weight:700;"
                        f"margin-bottom:4px;'>{zone_titles[zone_key]}</div>", unsafe_allow_html=True)
            if not people:
                st.markdown("<div style='text-align:center;color:#B7C0CC;font-size:11px;padding:8px;'>（無）</div>",
                            unsafe_allow_html=True)
            for person in people:
                st.markdown(f"<div style='text-align:center;background:#1F3A5F;color:#fff;border-radius:6px;"
                            f"padding:4px;font-size:12px;font-weight:700;margin-bottom:4px;'>{safe_text(person['Role'])}</div>",
                            unsafe_allow_html=True)
                st.markdown(person_card(person), unsafe_allow_html=True)
            if zone_key == "right" and partner_rows_preview:
                lines = []
                for r in partner_rows_preview:
                    item = safe_text(r.get("Title", ""), "") or safe_text(r.get("Expertise", ""), "")
                    company = safe_text(r.get("Name", ""), "[廠商待補]")
                    line = f"{item}：{company}" if item else company
                    lines.append(f"<div style='font-size:9.5px;color:#333333;padding:2px 0;'>{line}</div>")
                st.markdown(
                    "<div style='border:1px solid #D8DFE8;border-radius:8px;overflow:hidden;margin-top:6px;'>"
                    "<div style='background:#2C3E50;color:#fff;text-align:center;font-weight:700;"
                    "font-size:11px;padding:4px;'>協力廠商</div>"
                    "<div style='background:#F4F7F9;padding:6px 10px;'>"
                    + "".join(lines) + "</div></div>", unsafe_allow_html=True)

    if middles:
        mc1, mc2, mc3 = st.columns([1, 1, 1])
        with mc2:
            st.markdown(f"<div style='text-align:center;background:#1F3A5F;color:#fff;border-radius:6px;"
                        f"padding:4px;font-size:12px;font-weight:700;margin:10px 0 4px;'>{safe_text(middles[0]['Role'])}</div>",
                        unsafe_allow_html=True)
            st.markdown(person_card(middles[0]), unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("##### 工作組別")

    group_names = []
    for g in df[df["Layer"].isin(["GroupLeader", "GroupMember"])]["GroupName"]:
        g = g.strip()
        if g and g not in group_names:
            group_names.append(g)

    n_cols = min(4, max(len(group_names), 1))
    group_cols = st.columns(n_cols)
    for idx, gname in enumerate(group_names):
        col = group_cols[idx % n_cols]
        with col:
            st.markdown(f"<div style='background:#4A4A4A;color:#fff;text-align:center;border-radius:6px;"
                        f"padding:5px;font-size:12px;font-weight:700;margin-top:8px;'>{gname}</div>",
                        unsafe_allow_html=True)
            leader_row = df[(df["Layer"] == "GroupLeader") & (df["GroupName"] == gname)]
            if len(leader_row):
                st.markdown("<div style='font-size:10px;color:#1F3A5F;font-weight:700;margin-top:4px;'>組長</div>", unsafe_allow_html=True)
                st.markdown(person_card(leader_row.to_dict("records")[0]), unsafe_allow_html=True)
            members = df[(df["Layer"] == "GroupMember") & (df["GroupName"] == gname)]
            if len(members):
                st.markdown("<div style='font-size:10px;color:#1F3A5F;font-weight:700;'>組員</div>", unsafe_allow_html=True)
                for _, m in members.iterrows():
                    badges = parse_badges(m["Badges"])
                    st.markdown(f"<div style='display:flex;justify-content:space-between;background:#fff;"
                                f"border:1px solid #D8DFE8;border-radius:6px;padding:3px 6px;margin-bottom:4px;font-size:11px;'>"
                                f"<span>{safe_text(m['Name'])} {safe_text(m['Title'],'')}</span>"
                                f"<span>{badge_html(badges)}</span></div>", unsafe_allow_html=True)

    st.markdown("---")
    stats = st.session_state.stats
    fcols = st.columns(4)
    for i, key in enumerate(["avg_years", "masters", "technicians", "groups"]):
        val = stats.get(key, "")
        parts = val.split(" ", 1)
        num = parts[0] if parts else ""
        lbl = parts[1] if len(parts) > 1 else ""
        with fcols[i]:
            st.markdown(f"<div style='text-align:center;background:#EEF1F6;border-radius:8px;padding:10px;'>"
                        f"<div style='font-size:24px;font-weight:800;color:#1F3A5F;'>{num}</div>"
                        f"<div style='font-size:12px;color:#666;'>{lbl}</div></div>", unsafe_allow_html=True)

# ---- Tab 2: 人力配置表 ----
with tab2:
    rows = build_staffing_table_rows(df)
    preview_df = pd.DataFrame(rows)
    st.dataframe(preview_df, use_container_width=True, hide_index=True)

    st.markdown("")
    try:
        pptx2 = build_staffing_table_pptx(df)
        st.download_button("⬇️ 另外匯出 Output 2：人力配置表 (.pptx 版本)", data=pptx2,
                            file_name="主要工作人力配置表.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        st.error(f"Output 2 (pptx) 產生失敗：{e}")

# ---- Tab 3: 人員經歷敘述（含 Gemini AI 智慧潤飾） ----
with tab3:
    top_l, top_r = st.columns([3, 2])
    with top_l:
        st.caption("僅呈現組長以上同仁（Layer: Top / SubTop / Advisor / Middle / GroupLeader）之經歷敘述。")
    with top_r:
        refine_clicked = st.button(
            "✨ 一鍵使用 Gemini 依據擬任工作潤飾履歷",
            use_container_width=True,
            help="依據 JobDescription（擬任工作）潤飾 BioNarrative（履歷）；"
                 "若所選模型 404 / Model Not Found，將自動切換備援模型。",
        )

    if refine_clicked:
        effective_key = get_effective_gemini_key()
        if not effective_key:
            st.warning("請先輸入 Gemini API Key")
        elif gemini_genai is None:
            st.error("尚未安裝 google-genai 套件，請確認 requirements.txt 是否包含 google-genai 並重新部署")
        else:
            model_candidates = build_gemini_model_candidates(
                st.session_state.gemini_model_choice, st.session_state.gemini_custom_model,
            )
            if not model_candidates:
                st.warning("請選擇 Gemini 模型，或於「自訂模型」輸入 Model ID")
            else:
                target_mask = st.session_state.df["Layer"].isin(SENIOR_LAYERS)
                target_idx = st.session_state.df[target_mask].index.tolist()
                total = len(target_idx)

                if total == 0:
                    st.info("目前無組長以上人員資料可供潤飾")
                else:
                    progress = st.progress(0, text=f"AI 潤飾處理中... (0/{total})")
                    fail_count = 0
                    used_models = set()
                    for i, idx in enumerate(target_idx):
                        row = st.session_state.df.loc[idx]
                        try:
                            new_bio, used_model = refine_bio_with_gemini(
                                effective_key,
                                row["Name"], row["JobDescription"], row["BioNarrative"],
                                model_candidates,
                            )
                            st.session_state.df.loc[idx, "BioNarrative"] = new_bio
                            used_models.add(used_model)
                        except Exception as e:
                            fail_count += 1
                            st.warning(f"{safe_text(row['Name'])} 履歷潤飾失敗：{e}")
                        progress.progress((i + 1) / total, text=f"AI 潤飾處理中... ({i + 1}/{total})")
                    progress.empty()

                    model_note = f"（使用模型：{', '.join(sorted(used_models))}）" if used_models else ""
                    if fail_count == 0:
                        st.toast(f"履歷 AI 潤飾完成！{model_note}", icon="✨")
                    else:
                        st.toast(f"履歷 AI 潤飾完成（{total - fail_count}/{total} 筆成功）{model_note}", icon="⚠️")
                    st.rerun()

    df = st.session_state.df  # 潤飾後重新取用最新資料

    senior_df = df[df["Layer"].isin(SENIOR_LAYERS)]
    order = {"Top": 0, "SubTop": 1, "Advisor": 2, "Middle": 3, "GroupLeader": 4}
    senior_df = senior_df.copy()
    senior_df["_order"] = senior_df["Layer"].map(order).fillna(9)
    senior_df["_sub"] = senior_df.apply(lambda r: 0 if (r["Layer"] == "Top" and classify_leadership_zone(r) != "right") else 1, axis=1)
    senior_df = senior_df.sort_values(["_order", "_sub"], kind="stable")

    for _, r in senior_df.iterrows():
        c1, c2 = st.columns([1, 4])
        with c1:
            raw = safe_open_image(get_photo_bytes(st.session_state.photos, r["PhotoName"]))
            if raw:
                st.image(raw, width=140)
            else:
                st.markdown("<div style='width:140px;height:140px;background:#C9CFD8;border-radius:8px;"
                            "display:flex;align-items:center;justify-content:center;color:#555;'>照片待補</div>",
                            unsafe_allow_html=True)
        with c2:
            badges = parse_badges(r["Badges"])
            cert_names = [BADGE_MAP.get(b, b) for b in badges]
            tail = clean_separator_join([safe_text(r["Title"], "")] + cert_names, sep="／")
            header_line = f"{safe_text(r['Role'])}　{safe_text(r['Name'])}"
            if tail:
                header_line += f"　{tail}"
            st.markdown(f"**{header_line}**")
            st.write(safe_text(r["BioNarrative"]))
        st.markdown("---")
